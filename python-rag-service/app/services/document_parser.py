# 导入未来版本注解
from __future__ import annotations

# 导入各种文件解析所需的库
import csv
import io
import json
import logging
import uuid
from pathlib import Path
from typing import Callable

# 导入 HTML/XML 解析库（类似 Java 的 Jsoup）
from bs4 import BeautifulSoup

# 导入配置类
from app.core.config import Settings
# 导入文档切片模型
from app.models.schemas import DocumentChunk

# 创建日志记录器
logger = logging.getLogger(__name__)

# 尝试导入 PDF 解析库 pypdf
try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover
    PdfReader = None

# 尝试导入 Word 文档解析库 python-docx
try:
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover
    DocxDocument = None

# 尝试导入 Excel 解析库 openpyxl
try:
    from openpyxl import load_workbook
except ImportError:  # pragma: no cover
    load_workbook = None

# 尝试导入 PPT 解析库 python-pptx
try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None


# 自定义异常：文档解析错误
class DocumentParserError(Exception):
    pass


# 文档解析服务类：负责解析各种格式的文档并提取文本
# 类似 Spring 中处理文件上传和解析的服务
class DocumentParserService:
    # 支持的文件类型集合
    SUPPORTED_TYPES = {
        "pdf",
        "doc",
        "docx",
        "xls",
        "xlsx",
        "ppt",
        "pptx",
        "txt",
        "html",
        "xml",
        "md",
        "csv",
        "json",
    }

    def __init__(self, settings: Settings) -> None:
        # 保存配置对象
        self.settings = settings

    # 类方法：检查文件类型是否支持
    # filename: 文件名；返回布尔值
    @classmethod
    def is_supported(cls, filename: str | None) -> bool:
        # 如果文件名为空或不包含扩展名，则不支持
        if not filename or "." not in filename:
            return False
        # 获取文件扩展名并检查是否在支持列表中
        return filename.rsplit(".", 1)[1].lower() in cls.SUPPORTED_TYPES

    # 解析上传的文件：从字节流中提取文本并分块
    # filename: 文件名；content: 文件字节；content_type: MIME 类型
    def parse_upload(self, filename: str, content: bytes, content_type: str | None) -> list[DocumentChunk]:
        # 调用内部方法提取文本
        text = self._extract_text(filename, content)
        # 构建元数据信息
        metadata = {
            "source": filename,  # 来源文件名
            "fileSize": len(content),  # 文件大小（字节）
            "contentType": content_type or "application/octet-stream",  # MIME 类型
        }
        # 将文本分割成块并返回
        return self._split_text(text, metadata)

    # 解析文件路径：从本地文件读取并解析
    # path: 文件路径（字符串或 Path 对象）
    def parse_path(self, path: str | Path) -> list[DocumentChunk]:
        file_path = Path(path)
        # 读取文件字节内容
        content = file_path.read_bytes()
        # 提取文本
        text = self._extract_text(file_path.name, content)
        # 构建元数据
        metadata = {"source": file_path.name, "filePath": str(file_path)}
        return self._split_text(text, metadata)

    # 核心方法：根据文件类型选择对应的解析器提取文本
    def _extract_text(self, filename: str, content: bytes) -> str:
        # 获取文件扩展名
        ext = filename.rsplit(".", 1)[1].lower() if "." in filename else ""
        # 定义文件类型到解析函数的映射表（策略模式）
        parsers: dict[str, Callable[[bytes], str]] = {
            "txt": self._parse_text,  # 纯文本文件
            "md": self._parse_text,  # Markdown 文件
            "json": self._parse_json,  # JSON 文件
            "csv": self._parse_csv,  # CSV 文件
            "html": self._parse_html,  # HTML 文件
            "xml": self._parse_html,  # XML 文件
            "pdf": self._parse_pdf,  # PDF 文件
            "docx": self._parse_docx,  # Word 文档
            "xlsx": self._parse_xlsx,  # Excel 表格
            "pptx": self._parse_pptx,  # PowerPoint 演示文稿
            "doc": self._parse_legacy_binary,  # 旧版 Word（不支持）
            "xls": self._parse_legacy_binary,  # 旧版 Excel（不支持）
            "ppt": self._parse_legacy_binary,  # 旧版 PPT（不支持）
        }
        # 获取对应的解析函数
        parser = parsers.get(ext)
        if parser is None:
            raise DocumentParserError(f"不支持的文件类型：{ext}")
        # 调用解析函数并去除首尾空白
        text = parser(content).strip()
        if not text:
            raise DocumentParserError("文档解析后内容为空")
        return text

    # 解析纯文本文件：直接解码字节流
    def _parse_text(self, content: bytes) -> str:
        return content.decode("utf-8", errors="ignore")

    # 解析 JSON 文件：加载并重新格式化为美观的 JSON 字符串
    def _parse_json(self, content: bytes) -> str:
        # 解析 JSON 对象
        payload = json.loads(content.decode("utf-8", errors="ignore"))
        # 重新序列化为带缩进的格式，ensure_ascii=False 保留中文字符
        return json.dumps(payload, ensure_ascii=False, indent=2)

    # 解析 CSV 文件：读取每一行并拼接为字符串
    def _parse_csv(self, content: bytes) -> str:
        # 使用 csv 模块读取（自动处理引号和转义）
        reader = csv.reader(io.StringIO(content.decode("utf-8", errors="ignore")))
        # 将每行转换为 "单元格 1, 单元格 2, ..." 格式
        return "\n".join([", ".join(row) for row in reader])

    # 解析 HTML/XML 文件：使用 BeautifulSoup 提取纯文本
    def _parse_html(self, content: bytes) -> str:
        # 创建 BeautifulSoup 对象
        soup = BeautifulSoup(content.decode("utf-8", errors="ignore"), "html.parser")
        # 提取所有文本内容，用换行符分隔
        return soup.get_text(separator="\n")

    # 解析 PDF 文件：使用 pypdf 提取每页文本
    def _parse_pdf(self, content: bytes) -> str:
        if PdfReader is None:
            raise DocumentParserError("缺少 pypdf 依赖，无法解析 PDF")
        # 从字节流创建 PDF 阅读器
        reader = PdfReader(io.BytesIO(content))
        # 提取所有页面的文本并用换行符连接
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    # 解析 Word 文档（.docx）：提取所有段落文本
    def _parse_docx(self, content: bytes) -> str:
        if DocxDocument is None:
            raise DocumentParserError("缺少 python-docx 依赖，无法解析 DOCX")
        # 从字节流加载文档
        document = DocxDocument(io.BytesIO(content))
        # 提取所有段落的文本
        return "\n".join(paragraph.text for paragraph in document.paragraphs)

    # 解析 Excel 文件（.xlsx）：提取每个工作表的行列数据
    def _parse_xlsx(self, content: bytes) -> str:
        if load_workbook is None:
            raise DocumentParserError("缺少 openpyxl 依赖，无法解析 XLSX")
        # 从字节流加载工作簿（data_only=True 读取公式计算结果）
        workbook = load_workbook(io.BytesIO(content), data_only=True)
        rows: list[str] = []
        # 遍历所有工作表
        for sheet in workbook.worksheets:
            # 添加工作表名称标记
            rows.append(f"[Sheet] {sheet.title}")
            # 遍历所有行
            for row in sheet.iter_rows(values_only=True):
                # 过滤空单元格并转换为字符串
                cells = [str(cell) for cell in row if cell is not None]
                if cells:
                    # 用 " | " 分隔各列
                    rows.append(" | ".join(cells))
        return "\n".join(rows)

    # 解析 PowerPoint 文件（.pptx）：提取幻灯片文本
    def _parse_pptx(self, content: bytes) -> str:
        if Presentation is None:
            raise DocumentParserError("缺少 python-pptx 依赖，无法解析 PPTX")
        # 从字节流加载演示文稿
        presentation = Presentation(io.BytesIO(content))
        lines: list[str] = []
        # 遍历所有幻灯片
        for index, slide in enumerate(presentation.slides, start=1):
            # 添加幻灯片编号标记
            lines.append(f"[Slide {index}]")
            # 遍历幻灯片中的所有形状
            for shape in slide.shapes:
                # 如果形状包含文本则提取
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
        return "\n".join(lines)

    # 解析旧版二进制 Office 文件：抛出错误提示不支持
    def _parse_legacy_binary(self, _: bytes) -> str:
        raise DocumentParserError("旧版 Office 二进制格式暂不支持，请转换为 docx/xlsx/pptx 后上传")

    # 文本分块方法：将长文本分割成多个片段（RAG 的核心步骤）
    # text: 输入文本；metadata: 元数据信息
    def _split_text(self, text: str, metadata: dict[str, object]) -> list[DocumentChunk]:
        # 标准化空白字符：将连续空白替换为单个空格
        normalized = " ".join(text.split())
        # 获取配置的切片大小（至少 200 字符）
        chunk_size = max(self.settings.chunk_size, 200)
        # 获取重叠大小（不超过切片大小的一半）
        overlap = min(self.settings.chunk_overlap, chunk_size // 2)
        chunks: list[DocumentChunk] = []
        start = 0
        # 滑动窗口切分文本
        while start < len(normalized):
            # 计算结束位置
            end = min(len(normalized), start + chunk_size)
            # 提取切片内容并去除首尾空白
            content = normalized[start:end].strip()
            if content:
                # 复制元数据并添加切片索引
                chunk_metadata = dict(metadata)
                chunk_metadata["chunkIndex"] = len(chunks)
                # 创建文档切片对象（使用 UUID 作为唯一标识）
                chunks.append(
                    DocumentChunk(
                        id=str(uuid.uuid4()),
                        content=content,
                        metadata=chunk_metadata,
                    )
                )
            # 如果已到达文本末尾则退出
            if end >= len(normalized):
                break
            # 移动起始位置（减去重叠部分以保持上下文连贯）
            start = end - overlap
        # 记录日志
        logger.info("document split into %s chunks", len(chunks))
        return chunks
